from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def add_text_box(slide, text, left, top, width, height, font_size, font_bold=False, color="#ffffff", alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_bottom = tf.margin_top = tf.margin_left = tf.margin_right = 0
    
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = font_bold
        p.font.color.rgb = hex_to_rgb(color)
        p.font.name = font_name
        p.alignment = alignment
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    SLATE_900 = "#0f172a"
    SLATE_800 = "#1e293b"
    BLUE_500 = "#3b82f6"
    EMERALD_500 = "#10b981"
    PURPLE_500 = "#8b5cf6"
    AMBER_500 = "#f59e0b"
    PINK_500 = "#ec4899"
    RED_500 = "#ef4444"
    TEXT_MAIN = "#ffffff"
    TEXT_SUB = "#94a3b8"

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    # Decor
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(prs.slide_width - Inches(3)), int(-Inches(1)), Inches(4), Inches(4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
    c1.line.fill.background()
    
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(-Inches(1)), int(prs.slide_height - Inches(2)), Inches(3), Inches(3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = hex_to_rgb("#1e40af")
    c2.fill.transparency = 0.8
    c2.line.fill.background()

    add_text_box(slide, "FINAL YEAR PROJECT", Inches(1), Inches(1), Inches(3), Inches(0.5), 14, True, BLUE_500)
    add_text_box(slide, "AnalytixAI", Inches(1), Inches(1.5), Inches(10), Inches(1.5), 60, True, TEXT_MAIN)
    add_text_box(slide, "Automated Multi-Source Data Analysis System", Inches(1), Inches(2.7), Inches(10), Inches(0.8), 28, False, "#bfdbfe")
    add_text_box(slide, "An advanced platform leveraging Artificial Intelligence to automate data cleaning,\nprocessing, and visualization across diverse domains.", Inches(1), Inches(3.6), Inches(9), Inches(1), 16, False, "#cbd5e1")

    # Details Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5), Inches(11.33), Pt(1))
    line.fill.solid() 
    line.fill.fore_color.rgb = hex_to_rgb("#334155")
    line.line.fill.background()

    # Details Columns
    col_w = Inches(2.5)
    # Col 1
    add_text_box(slide, "PRESENTER", Inches(1), Inches(5.2), col_w, Inches(0.3), 12, True, TEXT_SUB)
    add_text_box(slide, "Rishabh Mourya", Inches(1), Inches(5.5), col_w, Inches(0.5), 18, True, TEXT_MAIN)
    # Col 2
    add_text_box(slide, "ROLL NO", Inches(4), Inches(5.2), col_w, Inches(0.3), 12, True, TEXT_SUB)
    add_text_box(slide, "25KCTYCS35\nGrp 26", Inches(4), Inches(5.5), col_w, Inches(0.8), 18, True, TEXT_MAIN)
    # Col 3
    add_text_box(slide, "INTERNAL MENTOR", Inches(7), Inches(5.2), col_w, Inches(0.3), 12, True, TEXT_SUB)
    add_text_box(slide, "Shalini M", Inches(7), Inches(5.5), col_w, Inches(0.5), 18, True, TEXT_MAIN)
    # Col 4
    add_text_box(slide, "DATE", Inches(10), Inches(5.2), col_w, Inches(0.3), 12, True, TEXT_SUB)
    add_text_box(slide, "21st Feb 2026", Inches(10), Inches(5.5), col_w, Inches(0.5), 18, True, TEXT_MAIN)

    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid() 
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    add_text_box(slide, "Problem Statement", Inches(1), Inches(0.5), Inches(8), Inches(1), 36, True, TEXT_MAIN)
    add_text_box(slide, "Identifying the challenges in modern data analysis workflows", Inches(1), Inches(1.2), Inches(8), Inches(0.5), 16, False, TEXT_SUB)

    def make_prob_card(slide, x, y, color, title, text):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
        box.line.color.rgb = hex_to_rgb(SLATE_800)
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.15), Inches(2.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = hex_to_rgb(color)
        accent.line.fill.background()

        add_text_box(slide, title, x + Inches(0.5), y + Inches(0.3), Inches(4), Inches(0.5), 18, True, color)
        add_text_box(slide, text, x + Inches(0.5), y + Inches(0.8), Inches(4.5), Inches(1.2), 14, False, "#cbd5e1")

    make_prob_card(slide, Inches(1), Inches(2), RED_500, "CORE ISSUE", "Manual data analysis is time-consuming and prone to human errors, leading to unreliable insights.")
    make_prob_card(slide, Inches(6.8), Inches(2), AMBER_500, "CURRENT GAPS", "Lack of unified tools for analyzing diverse data formats (Sales, Finance, HR) in one place.")
    make_prob_card(slide, Inches(1), Inches(4.5), BLUE_500, "PAIN POINT", "Existing analytics tools are often too complex, requiring specialized technical knowledge.")
    make_prob_card(slide, Inches(6.8), Inches(4.5), EMERALD_500, "NEED", "Urgent requirement for an automated, code-free solution to clean and visualize data instantly.")

    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid() 
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    add_text_box(slide, "Proposed Solution: AnalytixAI", Inches(1), Inches(0.5), Inches(9), Inches(1), 32, True, TEXT_MAIN)
    add_text_box(slide, "Automated Analytics", Inches(1), Inches(1.4), Inches(4), Inches(0.3), 14, True, BLUE_500)
    
    desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.9), Inches(5), Inches(1.4))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
    desc_box.line.color.rgb = hex_to_rgb(BLUE_500)
    desc_box.line.width = Pt(2)
    add_text_box(slide, "A full-stack platform that bridges the gap between raw data and actionable insights without writing code.", Inches(1.2), Inches(2.1), Inches(4.6), Inches(1.0), 14, False, "#e2e8f0")

    y_pos = 4.0
    for title, sub in [("Universal Ingestion", "CSV & Excel Support"), ("Multi-Domain", "Sales, Finance, HR"), ("Real-Time Processing", "Zero-latency results")]:
        add_text_box(slide, "• " + title, Inches(1), Inches(y_pos), Inches(5), Inches(0.4), 18, True, TEXT_MAIN)
        add_text_box(slide, "   " + sub, Inches(1), Inches(y_pos + 0.35), Inches(5), Inches(0.3), 14, False, TEXT_SUB)
        y_pos += 1.0

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(2), Pt(2), Inches(4.5))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb("#334155")
    line.line.fill.background()

    y_w = 2.0
    for num, text, color in [("1", "Data Upload", BLUE_500), ("2", "Auto-Processing", AMBER_500), ("3", "Visual Generation", PURPLE_500), ("4", "AI Insights", EMERALD_500)]:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.35), Inches(y_w), Inches(0.3), Inches(0.3))
        s.fill.solid() 
        s.fill.fore_color.rgb = hex_to_rgb(color)
        s.line.fill.background()
        add_text_box(slide, text, Inches(9), Inches(y_w - 0.1), Inches(3), Inches(0.5), 20, True, TEXT_MAIN)
        y_w += 1.2

    # Slide 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid() 
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    add_text_box(slide, "System Architecture", Inches(0), Inches(0.5), prs.slide_width, Inches(1), 32, True, TEXT_MAIN, PP_ALIGN.CENTER)

    cx = int(prs.slide_width / 2)
    def make_layer(slide, y, w, h, color, title, tech):
        left = int(cx - (w/2))
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, h)
        s.fill.solid() 
        s.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
        s.line.color.rgb = hex_to_rgb(color)
        s.line.width = Pt(3)
        add_text_box(slide, title, left + Inches(0.2), y + Inches(0.2), w, Inches(0.4), 16, True, TEXT_MAIN)
        add_text_box(slide, tech, left + Inches(0.2), y + h - Inches(0.5), w, Inches(0.4), 12, False, color)

    make_layer(slide, Inches(1.5), Inches(8), Inches(1.0), BLUE_500, "Frontend Layer", "HTML5 • CSS3 • JavaScript")
    # Arrow
    ar = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, int(cx - Inches(0.1)), int(Inches(2.55)), Inches(0.2), Inches(0.2))
    ar.fill.solid()
    ar.fill.fore_color.rgb = hex_to_rgb("#475569")
    ar.line.color.rgb = hex_to_rgb("#475569")

    make_layer(slide, Inches(2.8), Inches(8), Inches(1.0), EMERALD_500, "Backend API Layer", "Python 3.9 • FastAPI")
    ar = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, int(cx - Inches(0.1)), int(Inches(3.85)), Inches(0.2), Inches(0.2))
    ar.fill.solid()
    ar.fill.fore_color.rgb = hex_to_rgb("#475569")
    ar.line.color.rgb = hex_to_rgb("#475569")

    # Split Layer
    ai_left = int(cx - Inches(4))
    ai_s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ai_left, Inches(4.1), Inches(3.9), Inches(1.2))
    ai_s.fill.solid() 
    ai_s.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
    ai_s.line.color.rgb = hex_to_rgb(PURPLE_500)
    ai_s.line.width = Pt(3)
    add_text_box(slide, "AI Engine", ai_left + Inches(0.2), Inches(4.3), Inches(3.5), Inches(0.5), 16, True, TEXT_MAIN)
    add_text_box(slide, "Google Gemini API", ai_left + Inches(0.2), Inches(4.8), Inches(3.5), Inches(0.4), 12, False, PURPLE_500)

    d_left = int(cx + Inches(0.1))
    d_s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, d_left, Inches(4.1), Inches(3.9), Inches(1.2))
    d_s.fill.solid() 
    d_s.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
    d_s.line.color.rgb = hex_to_rgb(AMBER_500)
    d_s.line.width = Pt(3)
    add_text_box(slide, "Data Engine", d_left + Inches(0.2), Inches(4.3), Inches(3.5), Inches(0.5), 16, True, TEXT_MAIN)
    add_text_box(slide, "Pandas • Seaborn", d_left + Inches(0.2), Inches(4.8), Inches(3.5), Inches(0.4), 12, False, AMBER_500)

    ar = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, int(cx - Inches(0.1)), int(Inches(5.35)), Inches(0.2), Inches(0.2))
    ar.fill.solid() 
    ar.fill.fore_color.rgb = hex_to_rgb("#475569")
    ar.line.color.rgb = hex_to_rgb("#475569")

    make_layer(slide, Inches(5.6), Inches(8), Inches(1.0), RED_500, "Persistence Layer", "MongoDB • SQLite")

    # Slide 5: Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid() 
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    add_text_box(slide, "Key Features", Inches(1), Inches(0.5), Inches(8), Inches(1), 32, True, TEXT_MAIN)

    def make_feature_card(slide, x, y, title, desc, color):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.4), Inches(2.2))
        s.fill.solid() 
        s.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
        s.line.color.rgb = hex_to_rgb("#334155")
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.4), Inches(0.1))
        strip.fill.solid() 
        strip.fill.fore_color.rgb = hex_to_rgb(color)
        strip.line.fill.background()
        
        add_text_box(slide, title, x + Inches(0.2), y + Inches(0.3), Inches(3), Inches(0.6), 16, True, TEXT_MAIN)
        add_text_box(slide, desc, x + Inches(0.2), y + Inches(0.8), Inches(3), Inches(1.3), 12, False, "#cbd5e1")

    make_feature_card(slide, Inches(0.9), Inches(2), "Multi-Domain", "Unified analysis for Student, Sales, Finance and Employee data.", BLUE_500)
    make_feature_card(slide, Inches(4.9), Inches(2), "Auto-Cleaning", "Intelligent preprocessing handles missing values and duplicates.", EMERALD_500)
    make_feature_card(slide, Inches(8.9), Inches(2), "Smart Visuals", "Instant histograms, pie charts, and trend lines.", PURPLE_500)
    make_feature_card(slide, Inches(2.9), Inches(4.8), "Robust Security", "Enterprise-grade protection with JWT Auth and privacy.", RED_500)
    make_feature_card(slide, Inches(6.9), Inches(4.8), "Reports & AI Chat", "Export PDF reports and chat with your dataset.", AMBER_500)

    # Slide 6: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    add_text_box(slide, "Technology Stack", Inches(0), Inches(0.5), prs.slide_width, Inches(1), 32, True, TEXT_MAIN, PP_ALIGN.CENTER)

    def make_tech_card(slide, x, y, title, items, color):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.4), Inches(2.2))
        s.fill.solid() 
        s.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
        s.line.color.rgb = hex_to_rgb(color)
        s.line.width = Pt(1.5)
        add_text_box(slide, title, x + Inches(0.2), y + Inches(0.2), Inches(3), Inches(0.5), 14, True, color)
        
        iy = y + Inches(0.8)
        for item in items:
            chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), iy, Inches(1.8), Inches(0.35))
            chip.fill.solid() 
            chip.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
            chip.line.color.rgb = hex_to_rgb("#334155")
            add_text_box(slide, item, x + Inches(0.2), iy, Inches(1.8), Inches(0.35), 10, False, TEXT_MAIN, PP_ALIGN.CENTER)
            iy += Inches(0.45)

    make_tech_card(slide, Inches(0.9), Inches(1.8), "Backend", ["Python 3.9", "FastAPI"], EMERALD_500)
    make_tech_card(slide, Inches(4.9), Inches(1.8), "Data Science", ["Pandas", "Scikit-Learn"], BLUE_500)
    make_tech_card(slide, Inches(8.9), Inches(1.8), "Visualization", ["Matplotlib", "Seaborn"], PURPLE_500)
    make_tech_card(slide, Inches(0.9), Inches(4.5), "Database", ["MongoDB", "SQLite"], AMBER_500)
    make_tech_card(slide, Inches(4.9), Inches(4.5), "Frontend", ["HTML/CSS", "JavaScript"], PINK_500)
    make_tech_card(slide, Inches(8.9), Inches(4.5), "Tools", ["VS Code", "Git"], TEXT_SUB)

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(SLATE_900)
    bg.line.color.rgb = hex_to_rgb(SLATE_900)

    add_text_box(slide, "Conclusion & Future Scope", Inches(1), Inches(0.5), Inches(10), Inches(1), 32, True, TEXT_MAIN)

    c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(4.5), Inches(4.5))
    c_box.fill.solid() 
    c_box.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
    c_box.line.color.rgb = hex_to_rgb("#334155")
    
    add_text_box(slide, "Project Outcome", Inches(1.2), Inches(2.5), Inches(4), Inches(0.5), 18, True, BLUE_500)
    add_text_box(slide, "The project successfully automates the end-to-end data analysis lifecycle, delivering a robust platform that bridges the gap between raw data and actionable intelligence.", Inches(1.2), Inches(3.2), Inches(4), Inches(2), 14, False, "#cbd5e1")
    
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.5), Inches(2.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = hex_to_rgb("#064e3b")
    badge.line.color.rgb = hex_to_rgb("#064e3b")
    add_text_box(slide, "Objective Achieved", Inches(2), Inches(5.5), Inches(2.5), Inches(0.5), 12, True, EMERALD_500, PP_ALIGN.CENTER)

    def make_future_card(slide, x, y, title, desc, color):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.0), Inches(2.0))
        s.fill.solid() 
        s.fill.fore_color.rgb = hex_to_rgb(SLATE_800)
        s.line.color.rgb = hex_to_rgb(SLATE_800)
        l = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(2.0))
        l.fill.solid() 
        l.fill.fore_color.rgb = hex_to_rgb(color)
        l.line.fill.background()
        
        add_text_box(slide, title, x + Inches(0.3), y + Inches(0.2), Inches(2.5), Inches(0.4), 14, True, "#ffffff")
        add_text_box(slide, desc, x + Inches(0.3), y + Inches(0.7), Inches(2.5), Inches(1.2), 11, False, "#94a3b8")

    make_future_card(slide, Inches(6.0), Inches(2), "Extended Sources", "SQL & JSON Support", AMBER_500)
    make_future_card(slide, Inches(9.5), Inches(2), "Collaboration", "Live Team Workspace", PURPLE_500)
    make_future_card(slide, Inches(6.0), Inches(4.5), "Predictive Analytics", "ML Forecasting", EMERALD_500)
    make_future_card(slide, Inches(9.5), Inches(4.5), "Mobile App", "iOS & Android", PINK_500)

    prs.save('AnalytixAI_Final_Presentation.pptx')
    print("Presentation saved to AnalytixAI_Final_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
