from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def create_presentation():
    prs = Presentation()

    # --- Design Helper Functions ---
    def apply_branding(slide, title_text):
        # Add a blue header bar
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), Inches(10), Inches(1.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue
        header_shape.line.color.rgb = RGBColor(0, 51, 102)

        # Get or create Title shape
        title = slide.shapes.title
        if not title:
            # Create a title text box if one doesn't exist (e.g. on blank slides)
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        
        # Position Title
        title.top = Inches(0.2)
        title.left = Inches(0.5)
        title.width = Inches(9)
        title.height = Inches(0.8)
        
        # Style Title
        tf = title.text_frame
        tf.clear() # Clear existing text
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = RGBColor(255, 255, 255) # White
        p.font.bold = True
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.LEFT

    def add_bullet_points(slide, points):
        # Try to find the content placeholder
        body = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body = shape
                break
        
        if not body:
             # Create a text box if no content placeholder (e.g. on blank slides)
            body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))

        tf = body.text_frame
        tf.clear()  # Clear default empty paragraph

        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(22)
            p.space_after = Pt(14)
            p.level = 0
            
            # Handle sub-points
            if point.startswith("   "):
                p.text = point.strip()
                p.level = 1
                p.font.size = Pt(18)

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue Background
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "AnalytixAI"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Automated Multi-Source Data Analysis System"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(200, 200, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Student Details
    details_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
    tf = details_box.text_frame
    p = tf.add_paragraph()
    p.text = "Rishabh Mourya | Roll No: 25KCTYCS35 | Group No: 26\nInternal Mentor: Shalini M\nDate: 21st February 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


    # --- Slide 2: Problem Statement ---
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    apply_branding(slide, "Problem Statement")
    add_bullet_points(slide, [
        "Existing data analysis is often manual and time-consuming.",
        "Lack of unified tools for analyzing diverse data formats (Sales, Finance, HR, Education).",
        "Complex tools require technical expertise, excluding non-technical users.",
        "Need for an automated solution to clean, process, and visualize data instantly."
    ])

    # --- Slide 3: Proposed Solution ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_branding(slide, "Proposed Solution: AnalytixAI")
    add_bullet_points(slide, [
        "An automated full-stack platform for multi-domain data analysis.",
        "Supported Domains: Sales, Finance, Employee, and Student data.",
        "Key Workflow:",
        "   1. User uploads data (CSV/Excel).",
        "   2. System auto-cleans and processes data.",
        "   3. Generates instant visualizations and statistical insights.",
        "   4. Provides AI-powered chat for data queries."
    ])

    # --- Slide 4: System Architecture (Diagram) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank Layout
    apply_branding(slide, "System Architecture")

    # Draw Diagram
    # 1. User
    user = slide.shapes.add_shape(MSO_SHAPE.SMILEY_FACE, Inches(1), Inches(3), Inches(1), Inches(1))
    user.fill.solid()
    user.fill.fore_color.rgb = RGBColor(255, 204, 0)
    user.text_frame.text = "User"
    
    # Arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.1), Inches(3.4), Inches(1), Inches(0.2))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(100, 100, 100)

    # 2. Frontend
    frontend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(2.5), Inches(1.5), Inches(2))
    frontend.fill.solid()
    frontend.fill.fore_color.rgb = RGBColor(100, 149, 237)
    frontend.text_frame.text = "Frontend\n(HTML/JS)"
    
    # Arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.8), Inches(3.4), Inches(1), Inches(0.2))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(100, 100, 100)
    
    # 3. Backend
    backend = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(2.5), Inches(1.5), Inches(2))
    backend.fill.solid()
    backend.fill.fore_color.rgb = RGBColor(60, 179, 113)
    backend.text_frame.text = "Backend\n(FastAPI)"
    
    # 4. Database (Cylinder)
    db = slide.shapes.add_shape(MSO_SHAPE.CAN, Inches(8.5), Inches(1.5), Inches(1.2), Inches(1.5))
    db.fill.solid()
    db.fill.fore_color.rgb = RGBColor(205, 92, 92)
    db.text_frame.text = "MongoDB"
    
    # 5. AI Service (Cloud)
    ai = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(8.5), Inches(4.5), Inches(1.5), Inches(1))
    ai.fill.solid()
    ai.fill.fore_color.rgb = RGBColor(147, 112, 219)
    ai.text_frame.text = "Gemini AI"

    # Connectors
    # Use Elbow connectors
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, backend.left + backend.width, backend.top + backend.height/2, db.left, db.top + db.height/2)
    line1.line.color.rgb = RGBColor(100, 100, 100)
    line1.line.width = Pt(2)
    
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, backend.left + backend.width, backend.top + backend.height/2, ai.left, ai.top + ai.height/2)
    line2.line.color.rgb = RGBColor(100, 100, 100)
    line2.line.width = Pt(2)


    # --- Slide 5: Key Features ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_branding(slide, "Key Features")
    add_bullet_points(slide, [
        "Multi-Domain Support: Generic analysis for different sectors.",
        "Automatic Data Cleaning: Handles missing values and duplicates.",
        "Interactive Visualizations: Histograms, Pie Charts, Line Graphs.",
        "Secure Authentication: JWT-based login system.",
        "PDF Reporting: Downloadable comprehensive analysis reports.",
        "Chat with Data: AI assistant to answer specific questions about the data."
    ])

    # --- Slide 6: Technology Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    apply_branding(slide, "Technology Stack")
    
    # Two Columns for Tech Stack
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    p = tf.add_paragraph()
    p.text = "Backend & Data:"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    backend_stack = ["Python 3.9+", "FastAPI", "Pandas & NumPy", "Matplotlib & Seaborn", "MongoDB"]
    for item in backend_stack:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)

    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = right_box.text_frame
    p = tf.add_paragraph()
    p.text = "Frontend & Tools:"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    frontend_stack = ["HTML5, CSS3, JS", "Glassmorphism UI", "Git & GitHub", "VS Code"]
    for item in frontend_stack:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(20)

    # --- Slide 7: Conclusion & Future Scope ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_branding(slide, "Conclusion & Future Scope")
    add_bullet_points(slide, [
        "Project successfully automates the data analysis lifecycle.",
        "Future Enhancements:",
        "   - Support for more data formats (JSON, SQL dumps).",
        "   - Real-time collaboration features.",
        "   - Advanced Machine Learning predictive models.",
        "   - Mobile Application integration."
    ])

    prs.save('AnalytixAI_Presentation_v2.pptx')
    print("Presentation generated successfully: AnalytixAI_Presentation_v2.pptx")

if __name__ == "__main__":
    create_presentation()
